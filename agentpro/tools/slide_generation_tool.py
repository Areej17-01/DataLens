from .base_tool import Tool
from pptx import Presentation
from typing import Any, Optional
from abc import ABC, abstractmethod
from pydantic import BaseModel
from pathlib import Path
import math
import requests
import json
import os

class SlideGenerationTool(Tool):
    name: str = "PowerPoint Slides Generator"
    description: str = "Creates a simple PowerPoint slide deck and saves it to the DataLens temp output folder."
    action_type: str = "ppt_generate"
    input_format: str = """
{
  "session_id": "string (optional session identifier)",
  "query": "string (ask or describe what you want on the slides)",
  "title": "string (presentation title)",
  "filename": "string (output filename, optional)",
  "slides": [
    {
      "title": "string (slide title)",
      "points": ["string (bullet point)", "string (bullet point)", ...]
    },
    ...
  ]
}"""
    
    def run(self, input_text: Any) -> str:
        """Generate a simple PowerPoint presentation and save it under DataLens temp output."""
        try:
            # Parse input if it's a string
            if isinstance(input_text, str):
                try:
                    input_data = json.loads(input_text)
                except json.JSONDecodeError:
                    input_data = {"query": input_text}
            else:
                input_data = input_text

            session_id = input_data.get("session_id", "default")
            query_text = input_data.get("query", "")
            title = input_data.get("title") or (query_text[:60].strip() if query_text else "Generated Slides")
            if not title:
                title = "Generated Slides"

            slides = input_data.get("slides")
            if not isinstance(slides, list):
                slides = []

            if not slides and query_text:
                # Fallback: create a simple slide from the user's query
                slides = [
                    {
                        "title": title,
                        "points": [query_text.strip()],
                    }
                ]

            filename = input_data.get("filename")
            if not filename:
                safe_title = "".join(
                    c if c.isalnum() or c in (" ", "_", "-") else "_" for c in title
                ).strip().replace(" ", "_")
                filename = f"{safe_title or 'slides'}.pptx"
            if not filename.lower().endswith(".pptx"):
                filename += ".pptx"

            repo_root = Path(__file__).resolve().parents[2]
            slide_dir = repo_root / "DataLens" / "backend" / "temp" / session_id / "output" / "slides"
            slide_dir.mkdir(parents=True, exist_ok=True)

            output_path = slide_dir / filename

            # Create presentation
            prs = Presentation()
            title_slide = prs.slides.add_slide(prs.slide_layouts[0])
            title_slide.shapes.title.text = title

            for slide_data in slides:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                slide.shapes.title.text = slide_data.get("title", "Untitled Slide")
                content = slide.placeholders[1]
                tf = content.text_frame
                for point in slide_data.get("points", []):
                    p = tf.add_paragraph()
                    p.text = str(point)

            prs.save(str(output_path))

            rel_url = f"/temp/{session_id}/output/slides/{output_path.name}"
            result = {
                "message": f"Created presentation '{title}' with {len(slides)} content slides.",
                "slide_paths": [rel_url],
                "slide_count": len(slides),
            }
            return json.dumps(result)

        except Exception as e:
            return json.dumps({"error": str(e)})
